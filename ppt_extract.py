import streamlit as st
from io import BytesIO
import os
import shutil
import zipfile
from pptx import Presentation # type: ignore

st.title("PPT Image Ripper")

def extract_images(pptx_path, output_path):

    pres = Presentation(pptx_path)
    if not os.path.exists(output_path):
        os.mkdir(output_path)
    image_count = 0
    alt_text_file = os.path.join(output_path, "alt_text.txt")
    with open(alt_text_file, 'w') as alt_file:

        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.shape_type==13:
                    image = shape.image
                    image_bytes = image.blob

                    image_filename = f'image_{image_count}.{image.ext}'
                    image_path = os.path.join(output_path, image_filename)

                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)

                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    if (alt_text == ""):
                        alt_text="No alt Text"
                    
                    alt_file.write(f'{image_filename}: {alt_text}\n')
                    image_count += 1
    alt_file.close()
    return f'Extracted {image_count} images from {pptx_path}'
    

def create_zip(dir_path):
    # Create a bytes buffer to hold the zip file
    buffer = BytesIO()

    # Create a zip file in the buffer
    with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for root, dirs, files in os.walk(dir_path):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, start=dir_path)
                zip_file.write(file_path, arcname)

    # Set buffer position to the beginning
    buffer.seek(0)
    return buffer

Output_Path = "Extracted_Images"
 
def main():
    submitted = False
    with st.form("File Upload"):
        st.write("Upload your desired powerpoint file")
        uploaded_files = st.file_uploader(label="Your PPTX file", type=['ppt','pptx'], accept_multiple_files=False)
        submitted = st.form_submit_button('Submit')
    if(submitted):
        
        st.write(uploaded_files.name)
        extracted = extract_images(uploaded_files.name, Output_Path)
        st.write(extracted)
        if extracted:
            zip_buffer = create_zip(Output_Path)
            st.download_button(
                label="Download Zip",
                data=zip_buffer,
                file_name="Extracted_Images.zip",
                mime="application/zip"
            )
            if os.path.exists(Output_Path):
                shutil.rmtree(Output_Path)



if __name__ == '__main__':
    main()